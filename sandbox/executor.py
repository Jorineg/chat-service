"""Sandbox executor — runs inside the gVisor container.

Single-socket bidirectional protocol with the chat service:
  CS → SB: {"type": "execute", "code": "..."}
  SB → CS: {"type": "db_query", "sql": "..."}     (during execution)
  CS → SB: {"type": "db_result", "rows": [...]}
  SB → CS: {"type": "result", "stdout": "...", ...} (execution done)
  CS → SB: {"type": "shutdown"}
"""

import hashlib
import json
import os
import socket
import struct
import sys
from io import StringIO
from pathlib import Path

from bridge_client import BridgeClient, fmt, _send, _recv

try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
except ImportError:
    pass

SOCKET_PATH = "/shared/sandbox.sock"
WORK_DIR = Path("/work")
MAX_OUTPUT = 50_000

_HEADER_FMT = "!I"
_HEADER_SIZE = struct.calcsize(_HEADER_FMT)


def _snapshot_work_dir() -> dict[str, float]:
    """Record mtime of all files under /work/ for change detection."""
    snap = {}
    for p in WORK_DIR.rglob("*"):
        if p.is_file() and not p.name.startswith("."):
            try:
                snap[str(p)] = p.stat().st_mtime
            except OSError:
                pass
    return snap


def _detect_new_files(before: dict[str, float]) -> list[dict]:
    """Find files created or modified since the snapshot."""
    new_files = []
    for p in WORK_DIR.rglob("*"):
        if p.is_file() and not p.name.startswith("."):
            path_str = str(p)
            try:
                mtime = p.stat().st_mtime
            except OSError:
                continue
            if path_str not in before or mtime > before[path_str]:
                try:
                    content = p.read_bytes()
                    sha = hashlib.sha256(content).hexdigest()
                    new_files.append({
                        "path": path_str,
                        "filename": p.name,
                        "size_bytes": len(content),
                        "content_hash": sha,
                    })
                except OSError:
                    pass
    return new_files


def _extract_text(path: Path, mime_type: str = "") -> str:
    """Extract text from a file based on its type."""
    suffix = path.suffix.lower()

    if suffix == ".pdf" or "pdf" in mime_type:
        import fitz
        doc = fitz.open(str(path))
        parts = [page.get_text() for page in doc]
        doc.close()
        return "\n\n".join(parts)

    if suffix == ".docx":
        import docx
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)

    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)

    if suffix in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join(str(c) if c is not None else "" for c in row))
            parts.append(f"=== {ws.title} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    if suffix == ".csv":
        return path.read_text(errors="replace")

    try:
        return path.read_text(errors="replace")
    except Exception:
        return f"(Cannot extract text from {suffix} file)"


def _build_env(bridge: BridgeClient) -> dict:
    """Build the execution environment with all sandbox functions."""
    pending_images: list[dict] = []

    def file_image(file_id_or_path=None, page=None, max_dim=None, **kwargs):
        """Queue an image for the LLM to see. Accepts a chat_files UUID or a local path."""
        ref = file_id_or_path or kwargs.get("id_or_path") or kwargs.get("ref") or kwargs.get("path")
        if not ref:
            raise ValueError("file_image requires a file ID or path as first argument")
        pending_images.append({"ref": str(ref), "page": page, "max_dim": max_dim})
        print(f"[Image queued: {ref}]")

    def file_text(file_id_or_path=None, **kwargs) -> str:
        """Extract text from a file. Works for PDF, docx, pptx, xlsx, csv, plain text."""
        ref = file_id_or_path or kwargs.get("id") or kwargs.get("id_or_path") or kwargs.get("path")
        if not ref:
            raise ValueError("file_text requires a file ID or /work/ path as first argument")

        if "/" in str(ref):
            local_path = Path(str(ref))
            if not local_path.exists():
                raise FileNotFoundError(f"File not found: {ref}")
            return _extract_text(local_path, "")

        info = bridge.file_info(ref)
        filename = info.get("filename", "")
        mime = info.get("mime_type", "")

        if info.get("extracted_text"):
            return info["extracted_text"]

        local_path = None
        for p in WORK_DIR.rglob("*"):
            if p.name == filename and p.is_file():
                local_path = p
                break
        if not local_path:
            raise FileNotFoundError(f"File '{filename}' not found in /work/")

        return _extract_text(local_path, mime)

    return {
        "db": bridge.db,
        "fmt": fmt,
        "file_info": bridge.file_info,
        "file_text": file_text,
        "file_image": file_image,
        "describe_image": bridge.describe_image,
        "download_file": bridge.download_file,
        "download_craft_file": bridge.download_craft_file,
        "download_url": bridge.download_url,
        "add_activity_entry": bridge.add_activity_entry,
        "update_project_status": bridge.update_project_status,
        "update_project_profile": bridge.update_project_profile,
        "_pending_images": pending_images,
    }


def _run_code(code: str, env: dict) -> dict:
    """Execute code with full Python access."""
    captured = StringIO()
    old_stdout = sys.stdout
    sys.stdout = captured
    try:
        exec(compile(code, "<code>", "exec"), env)

        last_val = None
        stripped = code.strip()
        try:
            last_val = eval(compile(stripped, "<code>", "eval"), env)
        except (SyntaxError, TypeError):
            lines = stripped.split("\n")
            if lines:
                last = lines[-1].strip()
                skip = ("#", "import", "from", "def", "class", "if", "for",
                        "while", "try", "with", "return", "raise", "assert",
                        "pass", "break", "continue", "print(", "print (")
                if last and not any(last.startswith(p) for p in skip):
                    if "=" not in last or last.count("=") == last.count("=="):
                        try:
                            last_val = eval(last, env)
                        except Exception:
                            pass

        output = captured.getvalue()
        if len(output) > MAX_OUTPUT:
            output = output[:MAX_OUTPUT] + f"\n... (truncated, {len(output)} total chars)"

        result = last_val
        if result is not None:
            try:
                if hasattr(result, "isoformat"):
                    result = result.isoformat()
                elif isinstance(result, (set, frozenset)):
                    result = list(result)
                json.dumps(result)
            except (TypeError, ValueError):
                result = str(result)

        return {"result": result, "output": output or None, "error": None}
    except Exception as e:
        output = captured.getvalue()
        return {"result": None, "output": output or None, "error": f"{type(e).__name__}: {e}"}
    finally:
        sys.stdout = old_stdout


def main():
    sock = socket.socket(socket.AF_UNIX, socket.SOCK_STREAM)
    sock.connect(SOCKET_PATH)

    bridge = BridgeClient(sock)
    env = _build_env(bridge)

    try:
        while True:
            msg = _recv(sock)

            if msg["type"] == "execute":
                snapshot = _snapshot_work_dir()
                result = _run_code(msg["code"], env)
                new_files = _detect_new_files(snapshot)
                pending_images = env.get("_pending_images", [])

                _send(sock, {
                    "type": "result",
                    "stdout": result["output"],
                    "result": result["result"],
                    "error": result["error"],
                    "pending_images": list(pending_images),
                    "new_files": new_files,
                })
                pending_images.clear()

            elif msg["type"] == "shutdown":
                break
    except ConnectionError:
        pass
    finally:
        sock.close()


if __name__ == "__main__":
    main()
